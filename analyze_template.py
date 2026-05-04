from pptx import Presentation

prs = Presentation('VHS_Data-Analytics-SQL-Python_final.pptx')

print("🔍 TEMPLATE ANALYSE\n")
print(f"📊 Slides gesamt: {len(prs.slides)}\n")

# Analyze slides
count = 0
for slide in prs.slides:
    if count >= 3:
        break
    print(f"Slide {count + 1}:")
    print(f"  Layout: {slide.slide_layout.name}")
    print(f"  Shapes: {len(slide.shapes)}")
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text = shape.text[:40].replace('\n', ' ')
            print(f"    - Text: {text}...")
    count += 1
    print()

